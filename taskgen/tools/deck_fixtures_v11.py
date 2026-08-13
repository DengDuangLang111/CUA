"""Deck fixtures for the v11.1 impress repair (RUNBOOK section 2).

Runs INSIDE the build container (ubuntu + libreoffice-impress + python-pptx),
mounted at /data; the container then converts each .pptx to .odp with a real
soffice. Slide contents mirror the text outlines the original v11 setups
tried (and structurally could not) convert with `soffice --convert-to odp`.
"""

DECKS = {
    "shelter_grant":       [("Riverbend Shelter Renovation Grant", [])],
    "induction_deck":      [("New Starter Induction", ["People Team"]),
                            ("Day One", ["Badges and IT accounts"])],
    "line_layout":         [("Assembly Line 3 Layout",
                             ["Station 1: Frame press", "Station 2: Weld cell",
                              "Station 3: Paint booth", "Station 4: Final QA"])],
    "volunteer_briefing":  [("Volunteer Briefing", ["Arrivals and sign-in", "Sorting the donation bins"]),
                            ("Warehouse Safety", ["Lift with your legs", "Keep aisles clear"]),
                            ("Delivery Routes", ["Van keys at the desk", "Call ahead for flats"]),
                            ("Closing the Shift", ["Count the crates", "Lock the roller door"])],
    "onboarding":          [("New Hire Onboarding", []), ("Prepared by People Operations", [])],
    "meridian_fund_review": [("Meridian Balanced Fund", []), ("Holdings overview", []), ("Fee schedule", [])],
    "deposition_prep":     [("Deposition Prep", ["Witness sequence", "Exhibit handling"])],
}

if __name__ == "__main__":
    from pptx import Presentation
    for name, slides in DECKS.items():
        prs = Presentation()
        for title, bullets in slides:
            layout = prs.slide_layouts[1] if bullets else prs.slide_layouts[5]
            sl = prs.slides.add_slide(layout)
            sl.shapes.title.text = title
            if bullets:
                body = sl.placeholders[1].text_frame
                for i, b in enumerate(bullets):
                    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
                    p.text = b
        prs.save("/data/%s.pptx" % name)
    print("pptx built:", len(DECKS))
